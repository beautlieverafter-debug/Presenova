"""PowerPoint extraction and text-only mutation helpers.

The writer always opens the uploaded package and edits existing text runs in
place. It never rebuilds slides or copies visual objects into a new deck.
"""

import logging
import os
import tempfile
import hashlib
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)


def _extract_paragraphs(text_frame) -> list[dict]:
    """Extract paragraph/run text without modifying the source object."""
    return [
        {
            "para_index": para_idx,
            "text": para.text,
            "runs": [
                {"run_index": run_idx, "text": run.text}
                for run_idx, run in enumerate(para.runs)
            ],
        }
        for para_idx, para in enumerate(text_frame.paragraphs)
    ]


def _extract_image_metadata(shape) -> dict:
    """Capture image identity, intrinsic size, placement, and relationship."""
    blob = shape.image.blob
    metadata = {
        'shape_index': None,
        'sha256': hashlib.sha256(blob).hexdigest(),
        'bytes': len(blob),
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
        'relationship_id': None,
        'relationship_target': None,
        'media_part': None,
        'intrinsic_width': None,
        'intrinsic_height': None,
        'format': getattr(shape.image, 'ext', None),
    }
    try:
        image = shape.image
        metadata['media_part'] = str(image.part.partname)
    except Exception:
        pass
    try:
        blip = shape._element.blipFill.blip
        relationship_id = blip.get(qn('r:embed')) or blip.get(qn('r:link'))
        metadata['relationship_id'] = relationship_id
        if relationship_id:
            metadata['relationship_target'] = str(shape.part.rels[relationship_id].target_ref)
    except Exception:
        pass
    try:
        from PIL import Image
        with Image.open(BytesIO(blob)) as image:
            metadata['intrinsic_width'], metadata['intrinsic_height'] = image.size
    except Exception:
        # Pillow is optional at runtime; package/media hashes remain authoritative.
        pass
    return metadata


def _extract_chart_metadata(shape, shape_index: int) -> dict:
    """Capture only an existing chart title text target; chart data is untouched."""
    metadata = {'shape_index': shape_index, 'title_paragraphs': []}
    try:
        chart = shape.chart
        if chart.has_title and chart.chart_title.has_text_frame:
            metadata['title_paragraphs'] = _extract_paragraphs(chart.chart_title.text_frame)
    except Exception:
        logger.debug('[ppt_processor] Could not inspect chart title.', exc_info=True)
    return metadata


def extract_slides(file_path: str) -> list[dict]:
    """Return rewrite-safe slide structure and all supported text targets."""
    prs = Presentation(file_path)
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        shape_list = list(slide.shapes)
        title_text = ''
        title_shape_index = None
        title_shape = slide.shapes.title
        if title_shape is not None and getattr(title_shape, 'has_text_frame', False):
            title_text = title_shape.text.strip()
            try:
                title_shape_index = shape_list.index(title_shape)
            except ValueError:
                title_shape_index = None

        if not title_text:
            for idx, shape in enumerate(shape_list):
                if getattr(shape, 'has_text_frame', False) and shape.text.strip():
                    title_shape_index = idx
                    title_text = shape.text.strip().splitlines()[0]
                    break

        textboxes = []
        tables_data = []
        charts_data = []
        images_data = []
        images = tables = charts = 0

        for shape_idx, shape in enumerate(shape_list):
            stype = shape.shape_type
            if stype == MSO_SHAPE_TYPE.PICTURE:
                images += 1
                image_metadata = _extract_image_metadata(shape)
                image_metadata['shape_index'] = shape_idx
                images_data.append(image_metadata)
                continue
            if stype == MSO_SHAPE_TYPE.TABLE:
                tables += 1
                cells = []
                for row_idx, row in enumerate(shape.table.rows):
                    for column_idx, cell in enumerate(row.cells):
                        cells.append({
                            'row_index': row_idx,
                            'column_index': column_idx,
                            'paragraphs': _extract_paragraphs(cell.text_frame),
                        })
                tables_data.append({'shape_index': shape_idx, 'cells': cells})
                continue
            if stype == MSO_SHAPE_TYPE.CHART:
                charts += 1
                charts_data.append(_extract_chart_metadata(shape, shape_idx))
                continue
            if not getattr(shape, 'has_text_frame', False):
                continue

            paragraphs = _extract_paragraphs(shape.text_frame)
            placeholder_type = None
            if getattr(shape, 'is_placeholder', False):
                try:
                    placeholder_type = str(shape.placeholder_format.type)
                except Exception:
                    placeholder_type = 'unknown'
            textboxes.append({
                'shape_index': shape_idx,
                'shape_name': shape.name,
                'is_placeholder': getattr(shape, 'is_placeholder', False),
                'placeholder_type': placeholder_type,
                'paragraphs': paragraphs,
                'full_text': '\n'.join(p['text'] for p in paragraphs),
            })

        slides_data.append({
            'slide_number': slide_idx,
            'title': title_text or f'Slide {slide_idx}',
            'title_shape_index': title_shape_index,
            'textboxes': textboxes,
            'tables_data': tables_data,
            'charts_data': charts_data,
            'images_data': images_data,
            'images': images,
            'tables': tables,
            'charts': charts,
            'shapes_total': len(shape_list),
        })

    logger.info("[ppt_processor] Extracted %s slides from '%s'.", len(slides_data), file_path)
    return slides_data


def get_presentation_metadata(file_path: str) -> dict:
    """Return high-level presentation metadata."""
    prs = Presentation(file_path)
    return {
        'slide_count': len(prs.slides),
        'width_emu': prs.slide_width,
        'height_emu': prs.slide_height,
        'width_inches': round(prs.slide_width / 914400, 2),
        'height_inches': round(prs.slide_height / 914400, 2),
    }


def _set_paragraph_text_preserving_format(paragraph, new_text: str) -> None:
    """Replace paragraph text while keeping all original run elements.

    Text is distributed proportionally over the original runs instead of
    collapsing them into one run. This keeps mixed font/color/bold segments,
    hyperlinks, language settings, and other run-level OOXML properties.
    """
    runs = list(paragraph.runs)
    if not runs:
        if new_text:
            paragraph.add_run().text = new_text
        return

    old_lengths = [len(run.text or '') for run in runs]
    old_total = sum(old_lengths)
    if old_total == 0:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ''
        return

    cursor = 0
    cumulative = 0
    for index, (run, old_length) in enumerate(zip(runs, old_lengths)):
        cumulative += old_length
        end = len(new_text) if index == len(runs) - 1 else round(
            len(new_text) * cumulative / old_total
        )
        run.text = new_text[cursor:end]
        cursor = end


def _replace_text_frame_text(text_frame, rewritten_paragraphs: list[str]) -> None:
    """Replace exactly one string for each existing paragraph."""
    original_paragraphs = list(text_frame.paragraphs)
    if len(original_paragraphs) != len(rewritten_paragraphs):
        raise ValueError(
            f'Paragraph count changed ({len(original_paragraphs)} -> '
            f'{len(rewritten_paragraphs)}).'
        )
    for paragraph, new_text in zip(original_paragraphs, rewritten_paragraphs):
        _set_paragraph_text_preserving_format(paragraph, str(new_text))


def _replace_textbox_text(shape, rewritten_paragraphs: list[str]) -> None:
    """Compatibility wrapper for normal shapes and table cells."""
    _replace_text_frame_text(shape.text_frame, rewritten_paragraphs)


def _replace_table_cells(shape, cells: list[dict]) -> None:
    """Replace table cell text without changing the table structure."""
    rows = shape.table.rows
    for cell_data in cells:
        row_index = cell_data.get('row_index')
        column_index = cell_data.get('column_index')
        if not isinstance(row_index, int) or not isinstance(column_index, int):
            raise ValueError('Table cell coordinates must be integers.')
        if row_index < 0 or row_index >= len(rows):
            raise ValueError('Table row index is out of range.')
        row = rows[row_index]
        if column_index < 0 or column_index >= len(row.cells):
            raise ValueError('Table column index is out of range.')
        paragraphs = cell_data.get('paragraphs')
        if not isinstance(paragraphs, list) or not all(isinstance(p, str) for p in paragraphs):
            raise ValueError('Table cell paragraphs must be a list of strings.')
        _replace_text_frame_text(row.cells[column_index].text_frame, paragraphs)


def _replace_chart_title(shape, paragraphs: list[str]) -> None:
    """Update an existing chart title only; never touch chart data/format."""
    chart = shape.chart
    if not chart.has_title or not chart.chart_title.has_text_frame:
        raise ValueError('Chart title does not exist or is not editable text.')
    _replace_text_frame_text(chart.chart_title.text_frame, paragraphs)


def update_presentation_text(
    input_path: str,
    rewritten_slides: list[dict],
    original_slides: list[dict],
    output_path: str,
) -> None:
    """Apply only validated text changes and atomically publish a PPTX."""
    del original_slides  # kept in the public signature for compatibility
    prs = Presentation(input_path)
    original_prs = Presentation(input_path)
    slide_list = list(prs.slides)

    for rewritten_slide in rewritten_slides:
        slide_number = rewritten_slide.get('slide_number')
        if not isinstance(slide_number, int) or not 1 <= slide_number <= len(slide_list):
            raise ValueError(f'Invalid slide number: {slide_number!r}.')
        slide = slide_list[slide_number - 1]
        shapes = list(slide.shapes)

        for textbox in rewritten_slide.get('textboxes', []):
            shape_index = textbox.get('shape_index')
            if not isinstance(shape_index, int) or not 0 <= shape_index < len(shapes):
                raise ValueError(f'Invalid shape index on slide {slide_number}.')
            shape = shapes[shape_index]
            if not getattr(shape, 'has_text_frame', False):
                raise ValueError(f'Shape {shape_index} on slide {slide_number} has no text frame.')
            paragraphs = textbox.get('paragraphs')
            if not isinstance(paragraphs, list) or not all(isinstance(p, str) for p in paragraphs):
                raise ValueError(f'Invalid textbox paragraphs on slide {slide_number}.')
            _replace_textbox_text(shape, paragraphs)

        for table in rewritten_slide.get('tables', []):
            shape_index = table.get('shape_index')
            if not isinstance(shape_index, int) or not 0 <= shape_index < len(shapes):
                raise ValueError(f'Invalid table shape index on slide {slide_number}.')
            shape = shapes[shape_index]
            if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                raise ValueError(f'Shape {shape_index} on slide {slide_number} is not a table.')
            _replace_table_cells(shape, table.get('cells', []))

        for chart in rewritten_slide.get('charts', []):
            shape_index = chart.get('shape_index')
            if not isinstance(shape_index, int) or not 0 <= shape_index < len(shapes):
                raise ValueError(f'Invalid chart shape index on slide {slide_number}.')
            shape = shapes[shape_index]
            if shape.shape_type != MSO_SHAPE_TYPE.CHART:
                raise ValueError(f'Shape {shape_index} on slide {slide_number} is not a chart.')
            paragraphs = chart.get('title_paragraphs', [])
            if paragraphs:
                _replace_chart_title(shape, paragraphs)

    _validate_integrity(original_prs, prs)

    output_parent = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(output_parent, exist_ok=True)
    temp_output = None
    try:
        with tempfile.NamedTemporaryFile(
            dir=output_parent, prefix='.presentation-', suffix='.pptx', delete=False
        ) as stream:
            temp_output = stream.name
        prs.save(temp_output)
        saved_prs = Presentation(temp_output)
        _validate_integrity(original_prs, saved_prs)
        os.replace(temp_output, output_path)
        temp_output = None
        logger.info("[ppt_processor] Saved improved presentation to '%s'.", output_path)
    finally:
        if temp_output and os.path.exists(temp_output):
            try:
                os.remove(temp_output)
            except OSError:
                logger.warning("[ppt_processor] Could not remove temporary output '%s'.", temp_output)


def _count_shape_types(slide) -> dict:
    counts = {'images': 0, 'tables': 0, 'charts': 0, 'placeholders': 0, 'total': 0}
    for shape in slide.shapes:
        counts['total'] += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            counts['images'] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            counts['tables'] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            counts['charts'] += 1
        if getattr(shape, 'is_placeholder', False):
            counts['placeholders'] += 1
    return counts


def _validate_integrity(original_prs: Presentation, rewritten_prs: Presentation) -> None:
    """Check structure that text-only changes are not allowed to alter."""
    original_slides = list(original_prs.slides)
    rewritten_slides = list(rewritten_prs.slides)
    if len(original_slides) != len(rewritten_slides):
        raise ValueError(
            f'Slide count changed ({len(original_slides)} -> {len(rewritten_slides)}).'
        )
    if (original_prs.slide_width, original_prs.slide_height) != (
        rewritten_prs.slide_width, rewritten_prs.slide_height
    ):
        raise ValueError('Presentation dimensions changed.')

    for index, (original, rewritten) in enumerate(
        zip(original_slides, rewritten_slides), start=1
    ):
        if original.slide_layout.name != rewritten.slide_layout.name:
            raise ValueError(f'Slide layout changed on slide {index}.')
        if _count_shape_types(original) != _count_shape_types(rewritten):
            raise ValueError(f'Shape structure changed on slide {index}.')
        for original_shape, rewritten_shape in zip(original.shapes, rewritten.shapes):
            if (original_shape.shape_type, original_shape.name) != (
                rewritten_shape.shape_type, rewritten_shape.name
            ):
                raise ValueError(f'Shape identity changed on slide {index}.')

    logger.info('[ppt_processor] Structural integrity validation passed.')
