"""
Text Extractor Service
Extracts text from PPTX and PDF files for AI analysis.
"""

import logging
import os
import zipfile

logger = logging.getLogger(__name__)

ALLOWED_REWRITE_EXTENSIONS  = {'.pptx'}           # formats that produce an improved file
ALLOWED_ANALYSIS_EXTENSIONS = {'.pptx', '.pdf'}   # formats accepted for analysis


def is_allowed_for_rewrite(filename: str) -> bool:
    _, ext = os.path.splitext(filename or '')
    return ext.lower() in ALLOWED_REWRITE_EXTENSIONS


def is_allowed_for_analysis(filename: str) -> bool:
    _, ext = os.path.splitext(filename or '')
    return ext.lower() in ALLOWED_ANALYSIS_EXTENSIONS


def get_extension(filename: str) -> str:
    _, ext = os.path.splitext(filename or '')
    return ext.lower()


def validate_file_content(file_path: str, filename: str) -> None:
    """Validate the file signature and required container parts.

    Extension checks alone are not sufficient because a renamed executable or
    corrupt archive would otherwise reach python-pptx/PyMuPDF and become a 500.
    """
    ext = get_extension(filename)
    try:
        if ext == '.pptx':
            if not zipfile.is_zipfile(file_path):
                raise ValueError('The uploaded file is not a valid PowerPoint (.pptx) package.')
            with zipfile.ZipFile(file_path) as archive:
                names = set(archive.namelist())
                if 'ppt/presentation.xml' not in names or '[Content_Types].xml' not in names:
                    raise ValueError('The uploaded file is not a valid PowerPoint (.pptx) package.')
        elif ext == '.pdf':
            with open(file_path, 'rb') as stream:
                if stream.read(5) != b'%PDF-':
                    raise ValueError('The uploaded file is not a valid PDF document.')
        else:
            raise ValueError(f"Unsupported file type for validation: '{ext}'")
    except (OSError, zipfile.BadZipFile) as exc:
        raise ValueError('The uploaded file is corrupt or unreadable.') from exc


# ─────────────────────────────────────────────────────────────────────────────
# PPTX Extraction
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract all visible text from a PPTX file as a flat string.
    Used for quality analysis (not for rewriting — ppt_processor handles that).
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)
        lines = []

        def append_shape_text(shape, output_lines: list[str]) -> None:
            shape_type = getattr(shape, 'shape_type', None)
            if shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        output_lines.append(' | '.join(cells))
                return
            if getattr(shape, 'has_text_frame', False):
                text = shape.text_frame.text.strip()
                if text:
                    output_lines.append(text)
            # Text inside grouped shapes is useful for analysis. The rewriter
            # deliberately does not mutate grouped/SmartArt content because
            # python-pptx cannot safely round-trip every graphic object.
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    append_shape_text(child, output_lines)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- Slide {slide_idx} ---")
            for shape in slide.shapes:
                append_shape_text(shape, lines)

        return "\n".join(lines)

    except Exception as exc:
        logger.error(f"[text_extractor] PPTX extraction failed: {exc}")
        raise RuntimeError(f"Failed to extract text from PPTX: {exc}") from exc


# ─────────────────────────────────────────────────────────────────────────────
# PDF Extraction
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from a PDF using PyMuPDF (fitz).
    Returns a flat string with page delimiters.
    """
    try:
        lines = []
        try:
            import fitz  # PyMuPDF
            with fitz.open(file_path) as doc:
                pages = ((page_idx, page.get_text('text')) for page_idx, page in enumerate(doc, start=1))
                for page_idx, raw_text in pages:
                    text = raw_text.strip()
                    if text:
                        lines.append(f"--- Page {page_idx} ---")
                        lines.append(text)
        except ImportError:
            # pypdf is a lighter, pure-Python fallback already used elsewhere
            # in this project; it is sufficient for text-only analysis.
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            for page_idx, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or '').strip()
                if text:
                    lines.append(f"--- Page {page_idx} ---")
                    lines.append(text)
        result = "\n".join(lines)
        logger.info(f"[text_extractor] Extracted {len(result)} chars from PDF '{file_path}'.")
        return result

    except Exception as exc:
        logger.error(f"[text_extractor] PDF extraction failed: {exc}")
        raise RuntimeError(f"Failed to extract text from PDF: {exc}") from exc


# ─────────────────────────────────────────────────────────────────────────────
# Dispatcher
# ─────────────────────────────────────────────────────────────────────────────

def get_all_text_for_analysis(file_path: str, filename: str) -> str:
    """
    Extract all text from a file for quality analysis.
    Dispatches to the correct extractor based on file extension.
    """
    ext = get_extension(filename)
    if ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.pdf':
        return extract_text_from_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file type for text extraction: '{ext}'")
