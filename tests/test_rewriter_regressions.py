"""Regression coverage for preservation, validation, and PDF analysis paths."""

import os
import tempfile
import unittest

from pptx import Presentation
from pptx.dml.color import RGBColor

from main import create_app
from services.gemini_service import _fallback_local_rewrite, validate_rewritten_slides
from services.ppt_processor import extract_slides, update_presentation_text
from services.text_extractor import validate_file_content
from tools.compare_pptx import compare_pptx_files


class RewriterRegressionTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.app = create_app()
        cls.app.config.update(TESTING=True, MAX_CONTENT_LENGTH=50 * 1024 * 1024)

    def _make_rich_pptx(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(500000, 500000, 5000000, 1000000)
        paragraph = textbox.text_frame.paragraphs[0]
        first = paragraph.add_run()
        first.text = 'we dont '
        first.font.bold = True
        first.font.color.rgb = RGBColor(255, 0, 0)
        second = paragraph.add_run()
        second.text = 'stop here'
        second.font.italic = True
        second.font.color.rgb = RGBColor(0, 0, 255)
        second.hyperlink.address = 'https://example.com'

        table = slide.shapes.add_table(2, 2, 500000, 1800000, 5000000, 1500000).table
        table.cell(0, 0).text = 'we dont'
        table.cell(0, 1).text = 'continue'
        table.cell(1, 0).text = 'safe'
        table.cell(1, 1).text = 'content'

        handle = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        handle.close()
        presentation.save(handle.name)
        return handle.name

    def test_mixed_run_formatting_hyperlink_and_table_survive(self):
        original_path = self._make_rich_pptx()
        output_path = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx').name
        try:
            slides = extract_slides(original_path)
            rewritten = _fallback_local_rewrite(slides)
            update_presentation_text(original_path, rewritten, slides, output_path)

            saved = Presentation(output_path)
            paragraph = saved.slides[0].shapes[0].text_frame.paragraphs[0]
            self.assertEqual(len(paragraph.runs), 2)
            self.assertTrue(paragraph.runs[0].font.bold)
            self.assertEqual(paragraph.runs[0].font.color.rgb, RGBColor(255, 0, 0))
            self.assertTrue(paragraph.runs[1].font.italic)
            self.assertEqual(paragraph.runs[1].font.color.rgb, RGBColor(0, 0, 255))
            self.assertEqual(paragraph.runs[1].hyperlink.address, 'https://example.com')
            self.assertEqual(len(saved.slides[0].shapes[1].table.rows), 2)
            self.assertEqual(len(saved.slides[0].shapes[1].table.columns), 2)

            report = compare_pptx_files(original_path, output_path)
            self.assertTrue(report['visually_identical_structure'])
            self.assertGreaterEqual(len(report['text_changes_detected']), 1)
        finally:
            for path in (original_path, output_path):
                if os.path.exists(path):
                    os.remove(path)

    def test_corrupt_pptx_is_rejected_before_python_pptx(self):
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as handle:
            handle.write(b'not a zip file')
            path = handle.name
        try:
            with self.assertRaises(ValueError):
                validate_file_content(path, 'broken.pptx')
            with self.app.test_client() as client:
                with open(path, 'rb') as stream:
                    response = client.post(
                        '/api/presentation-rewriter/submit',
                        data={'file': (stream, 'broken.pptx')},
                        content_type='multipart/form-data',
                    )
                self.assertEqual(response.status_code, 400)
        finally:
            os.remove(path)

    def test_invalid_ai_structure_is_rejected(self):
        with self.assertRaises(ValueError):
            validate_rewritten_slides(
                [{'slide_number': 1, 'textboxes': [], 'tables_data': []}],
                [{'slide_number': 1, 'textboxes': [{'shape_index': 99, 'paragraphs': ['bad']}]}],
            )

    def test_pdf_analysis_endpoint_accepts_valid_pdf(self):
        try:
            from pypdf import PdfWriter
        except ImportError:
            self.skipTest('pypdf is not installed')
        pdf_path = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf').name
        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        with open(pdf_path, 'wb') as stream:
            writer.write(stream)
        try:
            with self.app.test_client() as client:
                with open(pdf_path, 'rb') as stream:
                    response = client.post(
                        '/api/presentation-rewriter/analyze',
                        data={'file': (stream, 'analysis.pdf')},
                        content_type='multipart/form-data',
                    )
                self.assertEqual(response.status_code, 200)
                self.assertIsNone(response.get_json()['slides_analysed'])
        finally:
            os.remove(pdf_path)


if __name__ == '__main__':
    unittest.main()
