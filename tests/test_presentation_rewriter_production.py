"""
Production Quality & Audit Test Suite
AI Presentation Rewriter — Presenova
"""

import os
import sys
import tempfile
import unittest
from unittest.mock import patch

sys.path.insert(0, os.path.abspath(os.path.dirname(os.path.dirname(__file__))))

from pptx import Presentation
from main import create_app

app = create_app()
from services.text_extractor import is_allowed_for_rewrite, is_allowed_for_analysis
from services.download_service import validate_upload_size, get_download_path, ensure_download_folder, generate_output_filename
from services.gemini_service import _clean_json_response, _fallback_local_rewrite, _polish_text_fallback
from services.ppt_processor import extract_slides, update_presentation_text
from services.language_tool_service import check_grammar, summarise_grammar_issues
from tools.compare_pptx import compare_pptx_files


class TestPresentationRewriterProduction(unittest.TestCase):

    def setUp(self):
        app.config['TESTING'] = True
        self.client = app.test_client()

        # Create temporary PPTX
        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "welcome to presenova presentation"
        slide1.placeholders[1].text = "we dont stop here and we can improve text"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Key Project Metrics"
        tf = slide2.placeholders[1].text_frame
        tf.text = "First point about efficiency"
        p2 = tf.add_paragraph()
        p2.text = "Second point about accuracy"

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(tmp.name)
        tmp.close()
        self.sample_pptx_file = tmp.name

    def tearDown(self):
        if os.path.exists(self.sample_pptx_file):
            try:
                os.remove(self.sample_pptx_file)
            except Exception:
                pass

    def test_file_extension_validation(self):
        self.assertTrue(is_allowed_for_rewrite("test.pptx"))
        self.assertTrue(is_allowed_for_rewrite("test.PPTX"))
        self.assertFalse(is_allowed_for_rewrite("test.pdf"))
        self.assertFalse(is_allowed_for_rewrite("test.docx"))
        self.assertFalse(is_allowed_for_rewrite("test.exe"))

        self.assertTrue(is_allowed_for_analysis("test.pptx"))
        self.assertTrue(is_allowed_for_analysis("test.pdf"))
        self.assertFalse(is_allowed_for_analysis("test.txt"))

    def test_upload_size_validation(self):
        validate_upload_size(10 * 1024 * 1024)
        with self.assertRaises(ValueError):
            validate_upload_size(51 * 1024 * 1024)

    def test_path_traversal_protection(self):
        filename = "../../etc/passwd"
        with self.assertRaises(FileNotFoundError):
            get_download_path(filename)

    def test_json_sanitization(self):
        raw_markdown = "```json\n{\"slides\": []}\n```"
        cleaned = _clean_json_response(raw_markdown)
        self.assertEqual(cleaned, "{\"slides\": []}")

        raw_preamble = "Here is the response:\n{\"slides\": []}\nHope this helps!"
        cleaned = _clean_json_response(raw_preamble)
        self.assertEqual(cleaned, "{\"slides\": []}")

    def test_local_polish_fallback(self):
        raw_text = "we dont stop here and it's great"
        polished = _polish_text_fallback(raw_text)
        # Complete sentences now get a trailing period (professional presentation style)
        self.assertEqual(polished, "We do not stop here and it is great.")

        slides = [{
            "slide_number": 1,
            "textboxes": [{
                "shape_index": 0,
                "paragraphs": [{"text": "we cant miss this opportunity"}]
            }]
        }]
        rewritten = _fallback_local_rewrite(slides)
        self.assertEqual(rewritten[0]["textboxes"][0]["paragraphs"][0], "We cannot miss this opportunity.")

    def test_grammar_check_handles_matches_without_rule_id(self):
        class FakeMatch:
            def __init__(self):
                self.message = 'Missing article'
                self.context = 'the project manager'
                self.offset = 0
                self.errorLength = 4
                self.replacements = ['the project manager']

        fake_tool = type('FakeTool', (), {'check': lambda self, text: [FakeMatch()]})()

        with patch('services.language_tool_service._get_tool', return_value=fake_tool):
            matches = check_grammar('This is a test sentence for grammar checking.')
            summary = summarise_grammar_issues(matches)

        self.assertEqual(len(matches), 1)
        self.assertEqual(matches[0]['rule_id'], None)
        self.assertIn('Missing article', summary)

    def test_pptx_extraction_and_preservation(self):
        slides_data = extract_slides(self.sample_pptx_file)
        self.assertEqual(len(slides_data), 2)
        self.assertIn("welcome", slides_data[0]["title"].lower())

        rewritten_slides = [{
            "slide_number": 1,
            "textboxes": [
                {"shape_index": 0, "paragraphs": ["Welcome to Presenova Presentation"]},
                {"shape_index": 1, "paragraphs": ["We do not stop here, and we can improve text"]}
            ]
        }, {
            "slide_number": 2,
            "textboxes": [
                {"shape_index": 0, "paragraphs": ["Key Project Metrics"]},
                {"shape_index": 1, "paragraphs": ["First point regarding operational efficiency", "Second point regarding data accuracy"]}
            ]
        }]

        output_folder = ensure_download_folder()
        out_name = generate_output_filename("test_preservation.pptx")
        out_path = os.path.join(output_folder, out_name)

        update_presentation_text(self.sample_pptx_file, rewritten_slides, slides_data, out_path)
        self.assertTrue(os.path.exists(out_path))

        report = compare_pptx_files(self.sample_pptx_file, out_path)
        self.assertTrue(report["slide_count_match"])
        self.assertTrue(report["dimensions_match"])
        self.assertTrue(report["visually_identical_structure"])

        if os.path.exists(out_path):
            os.remove(out_path)

    def test_api_submit_endpoint_success(self):
        with open(self.sample_pptx_file, 'rb') as f:
            response = self.client.post(
                '/api/presentation-rewriter/submit',
                data={'file': (f, 'sample_test.pptx')},
                content_type='multipart/form-data'
            )

        self.assertEqual(response.status_code, 200)
        json_data = response.get_json()
        self.assertTrue(json_data['success'])
        self.assertIn('download_url', json_data)
        self.assertEqual(json_data['slides_processed'], 2)
        self.assertIn('quality_scores', json_data)

    def test_api_submit_invalid_file(self):
        with tempfile.NamedTemporaryFile(suffix='.txt') as tmp:
            data = {'file': (tmp, 'invalid.txt')}
            response = self.client.post(
                '/api/presentation-rewriter/submit',
                data=data,
                content_type='multipart/form-data'
            )
        self.assertEqual(response.status_code, 400)
        json_data = response.get_json()
        self.assertFalse(json_data['success'])

    def test_api_download_endpoint(self):
        with open(self.sample_pptx_file, 'rb') as f:
            res = self.client.post(
                '/api/presentation-rewriter/submit',
                data={'file': (f, 'sample_test.pptx')},
                content_type='multipart/form-data'
            )
        out_filename = res.get_json()['output_filename']

        download_res = self.client.get(f'/api/presentation-rewriter/download/{out_filename}')
        self.assertEqual(download_res.status_code, 200)
        self.assertEqual(download_res.mimetype, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')

        missing_res = self.client.get('/api/presentation-rewriter/download/non_existent.pptx')
        self.assertEqual(missing_res.status_code, 404)


if __name__ == '__main__':
    unittest.main()
