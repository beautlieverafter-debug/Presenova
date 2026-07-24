"""Programmatic PPTX preservation comparison used by the rewrite pipeline."""

import hashlib
import logging
import os
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


def _image_hash(shape) -> str | None:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    try:
        return hashlib.sha256(shape.image.blob).hexdigest()
    except Exception:
        return None


def _shape_signature(shape) -> dict:
    signature = {
        'name': shape.name,
        'shape_type': str(shape.shape_type),
        'left': shape.left,
        'top': shape.top,
        'width': shape.width,
        'height': shape.height,
        'rotation': shape.rotation,
        'is_placeholder': bool(getattr(shape, 'is_placeholder', False)),
        'image_sha256': _image_hash(shape),
    }
    if signature['is_placeholder']:
        try:
            signature['placeholder_idx'] = shape.placeholder_format.idx
            signature['placeholder_type'] = str(shape.placeholder_format.type)
        except Exception:
            signature['placeholder_idx'] = None
            signature['placeholder_type'] = None
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        signature['table_dimensions'] = (len(shape.table.rows), len(shape.table.columns))
    return signature


def analyze_slide_structure(slide) -> dict:
    """Extract structural metrics, including geometry and binary image identity."""
    metrics = {
        'layout_name': slide.slide_layout.name,
        'total_shapes': len(slide.shapes),
        'pictures': 0,
        'tables': 0,
        'charts': 0,
        'placeholders': 0,
        'textboxes': 0,
        'shape_signatures': [],
        'text_content': [],
    }
    for shape in slide.shapes:
        metrics['shape_signatures'].append(_shape_signature(shape))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            metrics['pictures'] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            metrics['tables'] += 1
            for row in shape.table.rows:
                metrics['text_content'].append(' | '.join(cell.text for cell in row.cells))
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            metrics['charts'] += 1
            try:
                if shape.chart.has_title and shape.chart.chart_title.has_text_frame:
                    metrics['text_content'].append(shape.chart.chart_title.text_frame.text)
            except Exception:
                pass
        if getattr(shape, 'is_placeholder', False):
            metrics['placeholders'] += 1
        if getattr(shape, 'has_text_frame', False):
            metrics['textboxes'] += 1
            metrics['text_content'].append(shape.text_frame.text)
    return metrics


def _local_name(tag: str) -> str:
    return tag.rsplit('}', 1)[-1]


def _normalised_xml(blob: bytes, remove_text: bool = False) -> bytes:
    root = ET.fromstring(blob)
    if remove_text:
        for element in root.iter():
            if _local_name(element.tag) == 't':
                element.text = ''
    return ET.tostring(root, encoding='utf-8')


def _package_mismatches(original_path: str, improved_path: str) -> list[str]:
    """Compare OOXML parts after removing only text runs from slide XML."""
    ignored = {'docProps/core.xml', 'docProps/app.xml'}
    mismatches = []
    with zipfile.ZipFile(original_path) as original, zipfile.ZipFile(improved_path) as improved:
        original_names = set(original.namelist()) - ignored
        improved_names = set(improved.namelist()) - ignored
        for name in sorted(original_names | improved_names):
            if name not in original_names or name not in improved_names:
                mismatches.append(name)
                continue
            original_blob = original.read(name)
            improved_blob = improved.read(name)
            if name.endswith('.xml') or name.endswith('.rels') or name == '[Content_Types].xml':
                remove_text = (
                    (name.startswith('ppt/slides/slide') or name.startswith('ppt/charts/'))
                    and name.endswith('.xml')
                )
                try:
                    original_blob = _normalised_xml(original_blob, remove_text)
                    improved_blob = _normalised_xml(improved_blob, remove_text)
                except ET.ParseError:
                    pass
            if original_blob != improved_blob:
                mismatches.append(name)
    return mismatches


def compare_pptx_files(original_path: str, improved_path: str) -> dict:
    """Return a detailed structural, media, text, and package comparison."""
    if not os.path.isfile(original_path):
        raise FileNotFoundError(f'Original file not found: {original_path}')
    if not os.path.isfile(improved_path):
        raise FileNotFoundError(f'Improved file not found: {improved_path}')

    original_prs = Presentation(original_path)
    improved_prs = Presentation(improved_path)
    report = {
        'original_file': original_path,
        'improved_file': improved_path,
        'slide_count_orig': len(original_prs.slides),
        'slide_count_impr': len(improved_prs.slides),
        'slide_count_match': len(original_prs.slides) == len(improved_prs.slides),
        'dimensions_match': (original_prs.slide_width == improved_prs.slide_width and original_prs.slide_height == improved_prs.slide_height),
        'structural_mismatches': [],
        'text_changes_detected': [],
        'package_mismatches': [],
        'image_count_orig': 0,
        'image_count_impr': 0,
        'image_count_match': True,
        'image_hashes_match': True,
        'image_positions_match': True,
        'media_package_match': True,
        'relationships_match': True,
    }
    if not report['slide_count_match']:
        report['structural_mismatches'].append('Slide count changed.')
        report['visually_identical_structure'] = False
        return report
    if not report['dimensions_match']:
        report['structural_mismatches'].append('Presentation dimensions changed.')

    for index, (original_slide, improved_slide) in enumerate(
        zip(original_prs.slides, improved_prs.slides), start=1
    ):
        original_metrics = analyze_slide_structure(original_slide)
        improved_metrics = analyze_slide_structure(improved_slide)
        if original_metrics['layout_name'] != improved_metrics['layout_name']:
            report['structural_mismatches'].append(f'Slide {index}: layout changed.')
        for key in ('total_shapes', 'pictures', 'tables', 'charts', 'placeholders', 'textboxes'):
            if original_metrics[key] != improved_metrics[key]:
                report['structural_mismatches'].append(
                    f'Slide {index}: {key} changed ({original_metrics[key]} -> {improved_metrics[key]}).'
                )
        if original_metrics['shape_signatures'] != improved_metrics['shape_signatures']:
            report['structural_mismatches'].append(f'Slide {index}: shape geometry/media identity changed.')
        original_images = [item for item in original_metrics['shape_signatures'] if item['image_sha256']]
        improved_images = [item for item in improved_metrics['shape_signatures'] if item['image_sha256']]
        report['image_count_orig'] += len(original_images)
        report['image_count_impr'] += len(improved_images)
        if len(original_images) != len(improved_images):
            report['image_count_match'] = False
        if [item['image_sha256'] for item in original_images] != [item['image_sha256'] for item in improved_images]:
            report['image_hashes_match'] = False
        if [
            (item['left'], item['top'], item['width'], item['height'])
            for item in original_images
        ] != [
            (item['left'], item['top'], item['width'], item['height'])
            for item in improved_images
        ]:
            report['image_positions_match'] = False
        if original_metrics['text_content'] != improved_metrics['text_content']:
            report['text_changes_detected'].append({
                'slide_number': index,
                'original_text': original_metrics['text_content'],
                'improved_text': improved_metrics['text_content'],
            })

    report['package_mismatches'] = _package_mismatches(original_path, improved_path)
    report['media_package_match'] = not any(name.startswith('ppt/media/') for name in report['package_mismatches'])
    report['relationships_match'] = not any(name.endswith('.rels') for name in report['package_mismatches'])
    if not report['image_count_match'] or not report['image_hashes_match'] or not report['image_positions_match']:
        report['structural_mismatches'].append('Embedded image count, identity, or placement changed.')
    report['visually_identical_structure'] = not report['structural_mismatches'] and not report['package_mismatches']
    logger.info(
        '[compare_pptx] slides=%s text_changes=%s structural_mismatches=%s package_mismatches=%s',
        report['slide_count_orig'], len(report['text_changes_detected']),
        len(report['structural_mismatches']), len(report['package_mismatches'])
    )
    return report


def main() -> None:
    import argparse
    parser = argparse.ArgumentParser(description='Compare original and improved PPTX files.')
    parser.add_argument('original_file')
    parser.add_argument('improved_file')
    args = parser.parse_args()
    report = compare_pptx_files(args.original_file, args.improved_file)
    print(report)


if __name__ == '__main__':
    main()
